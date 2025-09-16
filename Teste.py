from pptx import Presentation
from pptx.util import Inches, Pt

# Cria uma apresentação vazia
prs = Presentation()

# Dados dos slides
slides_data = [
    {"title": "NARMNT – Normas Regulamentares de Material de Manutenção",
     "content": "Exército Brasileiro\n\nSugestão de imagem: Logotipo do Exército ou viatura/armamento militar."},

    {"title": "Introdução à NARMNT",
     "content": "- Diretrizes e procedimentos para manutenção de material militar.\n"
                "- Objetivo: assegurar disponibilidade e confiabilidade operacional.\n"
                "- Aplicável a todas as unidades do Exército.\n"
                "- Importância: gestão eficiente de recursos, redução de falhas e aumento da vida útil do equipamento.\n"
                "Sugestão de imagem: Esquema de ciclo de manutenção ou viatura sendo inspecionada."},

    {"title": "Tipos de Manutenção segundo a NARMNT",
     "content": "- Preventiva: inspeções periódicas antes da falha.\n"
                "  Ex.: troca de filtros de motores de viaturas.\n"
                "- Corretiva: realizada após falha para restaurar funcionalidade.\n"
                "  Ex.: substituição de peças defeituosas em metralhadoras.\n"
                "- Preditiva: baseada em monitoramento de indicadores de desgaste.\n"
                "  Ex.: análise de rolamentos em helicópteros.\n"
                "Sugestão de imagem: Diagrama comparativo dos tipos de manutenção."},

    {"title": "Responsabilidades na execução da manutenção",
     "content": "- Comandantes de Unidades: assegurar cumprimento da NARMNT.\n"
                "- Oficiais de Material: supervisionar e controlar a manutenção.\n"
                "- Técnicos e Mecânicos: executar procedimentos conforme manuais e normas de segurança.\n"
                "Sugestão de imagem: Fluxograma de responsabilidades ou foto de equipe de manutenção militar."},

    {"title": "Registros e Controle de Manutenção",
     "content": "- Todos os procedimentos devem ser registrados detalhadamente.\n"
                "- Informações obrigatórias: data/hora, tipo de manutenção, responsável, peças substituídas.\n"
                "- Ex.: diário de bordo de uma viatura.\n"
                "Sugestão de imagem: Foto de diário de manutenção ou planilha de controle."},

    {"title": "Segurança na manutenção",
     "content": "- Uso de ferramentas adequadas e EPI (Equipamentos de Proteção Individual).\n"
                "- Checklists para prevenir acidentes com sistemas hidráulicos, elétricos ou mecânicos.\n"
                "- Garantia de integridade do pessoal e do material.\n"
                "Sugestão de imagem: Técnicos usando EPI ou sinalização de segurança em oficina militar."},

    {"title": "Ciclo de Vida do Material",
     "content": "- Manutenção ao longo de toda a vida útil do equipamento.\n"
                "- Combinação de preventiva, corretiva e inspeções periódicas.\n"
                "- Exemplo: viatura com manutenção mensal, corretiva quando necessário e inspeção anual.\n"
                "Sugestão de imagem: Linha do tempo ilustrando o ciclo de vida do equipamento."},

    {"title": "Conclusão",
     "content": "- A NARMNT garante prontidão operacional do Exército.\n"
                "- Reduz falhas e aumenta confiabilidade do equipamento.\n"
                "- Define responsabilidades, tipos de manutenção, controle documental e normas de segurança.\n"
                "- Fundamental para eficiência e segurança nas unidades militares.\n"
                "Sugestão de imagem: Viatura militar pronta para missão ou equipe em ação."},
]

# Função para adicionar slide
def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Layout de título e conteúdo
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    text_box = slide.placeholders[1]
    text_box.text = content

# Adiciona os slides
for slide in slides_data:
    add_slide(prs, slide["title"], slide["content"])

# Salva o arquivo
file_path = "/mnt/data/NARMNT_Presentation.pptx"
prs.save(file_path)
file_path
